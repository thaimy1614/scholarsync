from flask import Flask, request, send_file, jsonify
import os
import comtypes.client
from pptx import Presentation
from gtts import gTTS
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
import unicodedata
import re
from flask_cors import CORS

app = Flask(__name__)
CORS(app, resources={r"/convert": {"origins": "http://localhost:3000"}})

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    slide_texts = []
    for slide in prs.slides:
        notes = slide.notes_slide.notes_text_frame.text.strip() if slide.notes_slide.notes_text_frame else ""
        if not notes:
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    text = re.sub(r'Page\s+\d+|Trang\s+\d+', '', text, flags=re.IGNORECASE)
                    unwanted_phrases = ['Slide Title', 'Chú thích', 'Notes', 'Footer'] 
                    for phrase in unwanted_phrases:
                        text = text.replace(phrase, '').strip()
                    slide_text += text + " "
            notes = slide_text.strip()
        notes = unicodedata.normalize("NFC", notes)
        notes = ' '.join(notes.split())
        if notes: 
            slide_texts.append(notes)
    return slide_texts

def generate_audio(text, output_path, lang="vi", tld="com"):
    print(f"Processing text: '{text}' for {output_path}") 
    temp_dir = os.path.dirname(output_path)
    os.makedirs(temp_dir, exist_ok=True)
    print(f"Temp directory created or exists: {temp_dir}") 
    sanitized_text = re.sub(r'[^\w\s.,!?]', '', text) 
    if sanitized_text:
        try:
            tts = gTTS(text=sanitized_text, lang=lang, tld=tld)
            tts.save(output_path)
        except Exception as e:
            print(f"gTTS error: {str(e)}")
            with open(output_path, "w") as f:
                pass 
    else:
        with open(output_path, "w") as f:
            pass

def slide_to_image(pptx_path, slide_index, output_path):
    slides_dir = os.path.dirname(output_path)
    if not os.path.exists(slides_dir):
        os.makedirs(slides_dir, exist_ok=True)
        test_file = os.path.join(slides_dir, "test.txt")
        try:
            with open(test_file, "w") as f:
                f.write("test")
            os.remove(test_file)
        except PermissionError:
            raise PermissionError(f"No write permission in {slides_dir}")
    
    max_retries = 3
    comtypes.CoInitialize()
    for attempt in range(max_retries):
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        try:
            presentation = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
            slide = presentation.Slides(slide_index + 1)
            page_setup = presentation.PageSetup
            width = int(page_setup.SlideWidth / 72 * 96)
            height = int(page_setup.SlideHeight / 72 * 96) 
            slide.Export(os.path.abspath(output_path), "PNG", width, height)
            break
        except comtypes.COMError as e:
            print(f"Attempt {attempt + 1} failed for slide {slide_index}: {str(e)}")
            time.sleep(1)
        finally:
            if 'presentation' in locals():
                presentation.Close()
        powerpoint.Quit()
    else:
        print(f"Failed to export slide {slide_index} after {max_retries} attempts. Please export slide_{slide_index}.png manually to {slides_dir}.")
        return False
    comtypes.CoUninitialize()
    return True

def create_video(pptx_path, output_video_path):
    base_dir = os.path.dirname(os.path.abspath(pptx_path))
    os.makedirs(os.path.join(base_dir, "temp"), exist_ok=True)
    slides_dir = os.path.join(base_dir, "slides")
    os.makedirs(slides_dir, exist_ok=True)
    
    slide_texts = extract_text_from_pptx(pptx_path)
    
    clips = []
    
    for i, text in enumerate(slide_texts):
        slide_image_path = os.path.join(slides_dir, f"slide_{i}.png")
        if not slide_to_image(pptx_path, i, slide_image_path):
            print(f"Skipping slide {i} due to export failure. Please export slide_{i}.png manually to {slides_dir}.")
            if os.path.exists(slide_image_path):
                print(f"Using existing {slide_image_path} as fallback.")
            else:
                continue
        
        audio_path = os.path.join("temp", f"slide_{i}_audio.mp3")
        generate_audio(text, audio_path)
        
        img_clip = ImageClip(slide_image_path)
        if os.path.getsize(audio_path) > 0:
            audio_clip = AudioFileClip(audio_path)
            img_clip = img_clip.set_duration(audio_clip.duration)
            img_clip = img_clip.set_audio(audio_clip)
        else:
            img_clip = img_clip.set_duration(2)
        clips.append(img_clip)
    
    if not clips:
        raise ValueError("No clips were created. All slide exports failed. Please export images manually to slides/ and retry.")
    
    final_clip = concatenate_videoclips(clips, method="compose")
    final_clip.write_videofile(output_video_path, codec="libx264", fps=24)
    
    for i in range(len(slide_texts)):
        audio_path = os.path.join("temp", f"slide_{i}_audio.mp3")
        if os.path.exists(audio_path):
            os.remove(audio_path)
    os.rmdir(os.path.join(base_dir, "temp"))
    for file in os.listdir(slides_dir):
        os.remove(os.path.join(slides_dir, file))
    os.rmdir(slides_dir)

@app.route('/convert', methods=['POST'])
def convert_pptx_to_video():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part in the request'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not file.filename.endswith('.pptx'):
        return jsonify({'error': 'File must be a .pptx file'}), 400
    
    upload_dir = 'uploads'
    os.makedirs(upload_dir, exist_ok=True)
    pptx_path = os.path.join(upload_dir, file.filename)
    file.save(pptx_path)
    
    output_video_path = os.path.join(upload_dir, 'output_video.mp4')
    
    try:
        create_video(pptx_path, output_video_path)
        
        return send_file(output_video_path, mimetype='video/mp4', as_attachment=True, download_name='converted_video.mp4')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    
    finally:
        if os.path.exists(pptx_path):
            os.remove(pptx_path)
        if os.path.exists(upload_dir) and not os.listdir(upload_dir):
            os.rmdir(upload_dir)

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)